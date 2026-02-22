"""Generate CBT Companion pitch presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
BLUE = RGBColor(0x41, 0x7B, 0xF5)
LIGHT_BLUE = RGBColor(0x64, 0xB5, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT = RGBColor(0x7C, 0x4D, 0xFF)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
RED_ACCENT = RGBColor(0xEF, 0x53, 0x50)
CARD_BG = RGBColor(0x27, 0x27, 0x45)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_slide(slide, items, left=1.2, top=2.0, width=11, size=20, spacing=0.55, color=WHITE):
    """Add bullet points one by one."""
    for i, item in enumerate(items):
        y = top + i * spacing
        add_text(slide, left, y, width, 0.5, item, size=size, color=color)


def add_card(slide, left, top, width, height, color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_divider(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(top),
                                    Inches(11.333), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()


# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0, 2.0, 13.333, 1.2, "CBT Companion", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0, 3.2, 13.333, 0.8, "AI-Powered Cognitive Behavioral Therapy using Beck's Protocol",
         size=24, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_divider(slide, 4.2)
add_text(slide, 0, 4.5, 13.333, 0.6, "How we implemented evidence-based therapy in code",
         size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ==================== SLIDE 2: WHAT IS BECK'S CBT ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "What is Beck's CBT?", size=36, color=BLUE, bold=True)
add_divider(slide, 1.1)

items = [
    "Developed by Dr. Aaron T. Beck (1960s) - the gold standard of talk therapy",
    "Core idea: It's not events that upset us, it's how we THINK about them",
    "Therapist helps patient identify distorted thoughts and replace them with balanced ones",
    "Uses structured sessions with standardized tools (BDI-II, Thought Records)",
    "Backed by 50+ years of clinical research",
]
add_bullet_slide(slide, [f"  {x}" for x in items], top=1.6, spacing=0.65, size=22)

# Add the CBT triangle
add_card(slide, 4, 5.2, 5.3, 1.8, CARD_BG)
add_text(slide, 4.3, 5.3, 4.8, 0.5, "The CBT Model (Beck's Cognitive Triangle)", size=16, color=ORANGE, bold=True)
add_text(slide, 4.3, 5.85, 4.8, 1.0, "Situation  ->  Automatic Thought  ->  Emotion  ->  Behavior\n"
         "          Change the thought  =  Change the feeling", size=16, color=LIGHT_GRAY)

# ==================== SLIDE 3: OUR IMPLEMENTATION OVERVIEW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "How We Built It - The Big Picture", size=36, color=BLUE, bold=True)
add_divider(slide, 1.1)

# Three cards
for i, (title, desc, col) in enumerate([
    ("32-State Protocol", "Full Beck session structure\nfrom assessment to closing", BLUE),
    ("3-Agent AI System", "Mirrors real therapist roles:\nquestioner, analyzer, treater", ACCENT),
    ("Severity Routing", "Adapts treatment path based\non patient's depression level", GREEN),
]):
    x = 1 + i * 4
    add_card(slide, x, 1.7, 3.5, 2.5, CARD_BG)
    add_text(slide, x + 0.2, 1.85, 3.1, 0.5, title, size=22, color=col, bold=True)
    add_text(slide, x + 0.2, 2.45, 3.1, 1.2, desc, size=17, color=LIGHT_GRAY)

add_text(slide, 0.8, 4.6, 11.5, 0.5, "Tech Stack", size=22, color=ORANGE, bold=True)
items = [
    "  Backend: Python / Flask API deployed on HuggingFace Spaces (Docker)",
    "  LLM: Groq API with LLaMA 3.3 70B - generates warm, therapeutic responses",
    "  Database: Turso (libSQL) - stores sessions, BDI scores, patient profiles",
    "  Auth: JWT-based user authentication",
]
add_bullet_slide(slide, items, top=5.2, spacing=0.5, size=18, color=LIGHT_GRAY)

# ==================== SLIDE 4: FULL SESSION FLOW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "Session Flow - Matches Beck's Structure", size=36, color=BLUE, bold=True)
add_divider(slide, 1.1)

phases = [
    ("PRE-SESSION", ORANGE, [
        "1. BDI-II Assessment  -  21-item depression screening (Beck et al., 1996)",
        "2. Bridge  -  Connect to previous session",
        "3. Homework Review  -  Check on assigned tasks",
        "4. Agenda Setting  -  Collaboratively decide today's focus",
        "5. Psychoeducation  -  Teach the CBT model (Session 1)",
    ]),
    ("CORE THERAPY", BLUE, [
        "6. Severity Routing  -  Choose treatment path by BDI score",
        "7. Cognitive Restructuring  -  Beck's 6 Socratic Questions (see next slide)",
        "   OR  Behavioral Activation  -  For severe depression (BDI >= 29)",
    ]),
    ("POST-SESSION", GREEN, [
        "8. Schema Check  -  Identify core beliefs (Session 4+)",
        "9. Thought Record (DRDT)  -  Generate Beck's worksheet",
        "10. Summary + Feedback  -  Patient states takeaways in own words",
    ]),
]

y = 1.4
for phase_name, color, items in phases:
    add_text(slide, 1, y, 3, 0.4, phase_name, size=20, color=color, bold=True)
    y += 0.45
    for item in items:
        add_text(slide, 1.5, y, 10.5, 0.4, item, size=16, color=LIGHT_GRAY)
        y += 0.38
    y += 0.2

add_text(slide, 1, y + 0.1, 11, 0.4, "Source: full_protocol.py - 32 states total, mirrors Beck (1979) session format",
         size=14, color=RGBColor(0x88, 0x88, 0x88))

# ==================== SLIDE 5: THE 6 SOCRATIC QUESTIONS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "Core: Beck's 6 Socratic Questions", size=36, color=BLUE, bold=True)
add_text(slide, 0.8, 1.0, 11, 0.4, "From Judith Beck's Cognitive Behavior Therapy: Basics and Beyond (2020)",
         size=16, color=LIGHT_GRAY)
add_divider(slide, 1.4)

questions = [
    ("Q1", "Examining the Evidence", "What evidence supports this thought? What evidence is against it?",
     "Q1_EVIDENCE_FOR / Q1_EVIDENCE_AGAINST"),
    ("Q2", "Alternative Explanation", "Is there another way to look at this situation?",
     "Q2_ALTERNATIVE"),
    ("Q3", "Decatastrophizing", "What's the worst / best / most realistic outcome?",
     "Q3_WORST / Q3_BEST / Q3_REALISTIC"),
    ("Q4", "Effect of Believing", "How does holding this thought affect you?",
     "Q4_EFFECT"),
    ("Q5", "Double Standard", "What would you tell a friend with this same thought?",
     "Q5_FRIEND"),
    ("Q6", "Action Planning", "What can you do about this situation?",
     "Q6_ACTION"),
]

y = 1.7
for q_num, name, desc, code_ref in questions:
    add_card(slide, 0.8, y, 11.7, 0.72, CARD_BG)
    add_text(slide, 1.0, y + 0.05, 0.6, 0.4, q_num, size=18, color=BLUE, bold=True)
    add_text(slide, 1.7, y + 0.05, 3.5, 0.4, name, size=18, color=WHITE, bold=True)
    add_text(slide, 5.3, y + 0.05, 4.5, 0.4, desc, size=16, color=LIGHT_GRAY)
    add_text(slide, 10.0, y + 0.05, 2.3, 0.4, code_ref, size=12, color=RGBColor(0x88, 0x88, 0x88))
    y += 0.82

add_text(slide, 0.8, y + 0.15, 11, 0.4,
         "Each question = a separate state in our state machine. User answers one at a time, just like real therapy.",
         size=16, color=ORANGE)

# ==================== SLIDE 6: 3-AGENT SYSTEM ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "3-Agent AI System", size=36, color=BLUE, bold=True)
add_text(slide, 0.8, 1.0, 11, 0.4, "Mirrors how a real CBT therapist thinks and acts",
         size=18, color=LIGHT_GRAY)
add_divider(slide, 1.45)

agents = [
    ("Agent 1", "Warm Questioner", BLUE,
     "Handles: VALIDATE through Q6_ACTION (13 states)",
     "Validates emotions, asks Beck's Socratic questions one\n"
     "at a time, captures belief ratings and emotions.\n"
     "Warm and empathic tone throughout."),
    ("Agent 2", "Clinical Summarizer", ACCENT,
     "Handles: SUMMARIZING (1 internal state)",
     "Analyzes all patient responses internally.\n"
     "Finds contradictions in the evidence.\n"
     "Patient never sees this - just like a real therapist's notes."),
    ("Agent 3", "Treatment Agent", GREEN,
     "Handles: DELIVER_REFRAME through COMPLETE (6 states)",
     "Delivers balanced thought using patient's OWN words.\n"
     "Re-rates belief and emotion (measures change).\n"
     "Creates homework / behavioral experiment."),
]

for i, (name, role, color, handles, desc) in enumerate(agents):
    x = 0.8 + i * 4.1
    add_card(slide, x, 1.8, 3.8, 4.5, CARD_BG)
    add_text(slide, x + 0.25, 1.95, 3.3, 0.4, name, size=24, color=color, bold=True)
    add_text(slide, x + 0.25, 2.45, 3.3, 0.4, role, size=18, color=WHITE, bold=True)
    add_text(slide, x + 0.25, 2.95, 3.3, 0.4, handles, size=13, color=RGBColor(0x88, 0x88, 0x88))
    add_text(slide, x + 0.25, 3.5, 3.3, 2.2, desc, size=15, color=LIGHT_GRAY)

add_text(slide, 0.8, 6.6, 11.5, 0.4, "Source: groq_client.py (Agents 1-3) + beck_agents.py (extended protocol agents)",
         size=14, color=RGBColor(0x88, 0x88, 0x88))

# ==================== SLIDE 7: BDI-II IMPLEMENTATION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "BDI-II: Beck's Own Assessment Tool", size=36, color=BLUE, bold=True)
add_divider(slide, 1.1)

add_text(slide, 0.8, 1.4, 11, 0.4, "We implement the full Beck Depression Inventory-II (Beck et al., 1996)",
         size=20, color=WHITE)

# Left side - items
add_card(slide, 0.8, 2.1, 5.5, 4.5, CARD_BG)
add_text(slide, 1.0, 2.2, 5, 0.4, "All 21 BDI-II Items", size=20, color=ORANGE, bold=True)
bdi_items = [
    "Sadness, Pessimism, Past Failure, Loss of Pleasure,",
    "Guilty Feelings, Punishment Feelings, Self-Dislike,",
    "Self-Criticalness, Suicidal Thoughts (CRISIS ITEM),",
    "Crying, Agitation, Loss of Interest, Indecisiveness,",
    "Worthlessness, Loss of Energy, Sleep Changes,",
    "Irritability, Appetite Changes, Concentration,",
    "Tiredness, Loss of Interest in Sex",
]
y = 2.7
for item in bdi_items:
    add_text(slide, 1.2, y, 4.8, 0.35, item, size=14, color=LIGHT_GRAY)
    y += 0.35

# Right side - scoring
add_card(slide, 6.8, 2.1, 5.7, 4.5, CARD_BG)
add_text(slide, 7.0, 2.2, 5.2, 0.4, "Scoring (Beck's exact cutoffs)", size=20, color=ORANGE, bold=True)

scores = [
    ("0 - 13", "Minimal depression", GREEN),
    ("14 - 19", "Mild depression", RGBColor(0xFF, 0xEB, 0x3B)),
    ("20 - 28", "Moderate depression", ORANGE),
    ("29 - 63", "Severe depression", RED_ACCENT),
]
y = 2.85
for score_range, label, color in scores:
    add_text(slide, 7.2, y, 2.0, 0.45, score_range, size=20, color=color, bold=True)
    add_text(slide, 9.3, y, 3.0, 0.45, label, size=18, color=LIGHT_GRAY)
    y += 0.6

add_text(slide, 7.2, y + 0.3, 5.0, 0.8,
         "Each item scored 0-3 by patient\nItem 9 (Suicidal Thoughts) triggers crisis protocol",
         size=15, color=RGBColor(0xAA, 0xAA, 0xAA))

# ==================== SLIDE 8: SEVERITY ROUTING ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "Smart Treatment Routing", size=36, color=BLUE, bold=True)
add_text(slide, 0.8, 1.0, 11, 0.4,
         "Based on Beck et al. (1979): Severely depressed patients need behavioral activation BEFORE cognitive work",
         size=17, color=LIGHT_GRAY)
add_divider(slide, 1.45)

# Three routing paths
paths = [
    ("SEVERE", "BDI >= 29", RED_ACCENT,
     "Behavioral Activation",
     "Too depressed for thought work.\nFocus on DOING, not thinking.\n\n"
     "1. Activity Monitoring\n2. Activity Scheduling\n3. Graded Task Assignment\n\n"
     "\"Activity precedes mood change\""),
    ("MILD / MODERATE", "BDI 14-28", BLUE,
     "Cognitive Restructuring",
     "Standard Beck protocol.\nThe 6 Socratic Questions.\n\n"
     "Identify distorted thought\n-> Examine evidence\n-> Generate balanced thought\n-> Measure improvement"),
    ("RECOVERED", "BDI < 14 (3+ sessions)", GREEN,
     "Relapse Prevention",
     "Patient is getting better!\nPrepare for independence.\n\n"
     "1. Celebrate BDI progress\n2. Identify warning signs\n3. Build coping plan\n4. Space out sessions"),
]

for i, (level, criteria, color, treatment, desc) in enumerate(paths):
    x = 0.6 + i * 4.2
    add_card(slide, x, 1.8, 3.9, 5.0, CARD_BG)
    add_text(slide, x + 0.25, 1.95, 3.4, 0.4, level, size=22, color=color, bold=True)
    add_text(slide, x + 0.25, 2.4, 3.4, 0.35, criteria, size=16, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text(slide, x + 0.25, 2.9, 3.4, 0.4, treatment, size=18, color=WHITE, bold=True)
    add_text(slide, x + 0.25, 3.4, 3.4, 3.0, desc, size=15, color=LIGHT_GRAY)

# ==================== SLIDE 9: ADVANCED BECK TECHNIQUES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "Advanced Beck Techniques We Implement", size=36, color=BLUE, bold=True)
add_divider(slide, 1.1)

techniques = [
    ("Downward Arrow Technique", "Beck (1995)",
     "Drills from surface thought to CORE BELIEF in 3-5 levels.\n"
     "\"If that were true, what would it mean about you?\"\n"
     "Identifies: Helplessness / Unlovability / Worthlessness",
     "schema_agent() in beck_agents.py", ACCENT),
    ("Daily Record of Dysfunctional Thoughts", "Beck (1979)",
     "Auto-generates Beck's DRDT worksheet after each session:\n"
     "Situation -> Thought -> Emotion -> Evidence -> Balanced Thought -> Outcome",
     "drdt_agent() in beck_agents.py", BLUE),
    ("Collaborative Empiricism", "Beck (1979)",
     "Patient and therapist work TOGETHER - not top-down.\n"
     "Agenda is set collaboratively. Feedback is explicitly requested.\n"
     "Patient restates takeaways in own words (guards against compliance).",
     "agenda_setting_agent(), feedback_agent()", GREEN),
    ("Behavioral Experiments", "Beck (1979)",
     "Each session ends with a small, testable action plan.\n"
     "Reviewed at start of next session as homework.\n"
     "Turns therapy insights into real-world evidence.",
     "ACTION_PLAN state in prompts.py", ORANGE),
]

y = 1.4
for name, source, desc, code, color in techniques:
    add_card(slide, 0.8, y, 11.7, 1.2, CARD_BG)
    add_text(slide, 1.1, y + 0.08, 5.0, 0.35, name, size=19, color=color, bold=True)
    add_text(slide, 6.2, y + 0.08, 2.5, 0.35, source, size=14, color=RGBColor(0x88, 0x88, 0x88))
    add_text(slide, 1.1, y + 0.45, 7.5, 0.7, desc, size=14, color=LIGHT_GRAY)
    add_text(slide, 8.8, y + 0.08, 3.5, 0.35, code, size=12, color=RGBColor(0x77, 0x77, 0x77))
    y += 1.35

# ==================== SLIDE 10: COMPLETE STATE MAP ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "All 32 States at a Glance", size=36, color=BLUE, bold=True)
add_divider(slide, 1.0)

# Pre-session column
add_text(slide, 0.5, 1.2, 3.0, 0.3, "PRE-SESSION (6)", size=16, color=ORANGE, bold=True)
pre = [
    ("BDI_ASSESSMENT", "Screen depression (21 items)"),
    ("BRIDGE", "Connect to last session"),
    ("HOMEWORK_REVIEW", "Check assigned tasks"),
    ("AGENDA_SETTING", "Pick today's focus together"),
    ("PSYCHOEDUCATION", "Teach CBT model (Session 1)"),
    ("SEVERITY_ROUTING", "Pick treatment path by score"),
]
y = 1.55
for s, desc in pre:
    add_card(slide, 0.5, y, 4.0, 0.32, CARD_BG)
    add_text(slide, 0.6, y + 0.03, 2.2, 0.25, s, size=11, color=ORANGE, bold=True)
    add_text(slide, 2.85, y + 0.03, 1.6, 0.25, desc, size=10, color=LIGHT_GRAY)
    y += 0.35

# Arrow + routing label
add_text(slide, 0.5, y + 0.05, 4.0, 0.35,
         "Routes by BDI:  >= 29 Severe  |  14-28 Moderate  |  < 14 Recovered",
         size=10, color=WHITE, bold=True)

# Cognitive column
add_text(slide, 5.0, 1.2, 4.0, 0.3, "COGNITIVE FLOW (20 states)", size=16, color=BLUE, bold=True)
cog = [
    ("VALIDATE", "Acknowledge the pain"),
    ("RATE_BELIEF", "How much you believe it? 0-100%"),
    ("CAPTURE_EMOTION", "What emotion comes up?"),
    ("RATE_EMOTION", "How intense? 0-100"),
    ("Q1: EVIDENCE FOR", "What supports this thought?"),
    ("Q1: EVIDENCE AGAINST", "Any exceptions?"),
    ("Q2: ALTERNATIVE", "Another explanation?"),
    ("Q3: WORST / BEST / REALISTIC", "Decatastrophize"),
    ("Q4: EFFECT", "How does believing it affect you?"),
    ("Q5: FRIEND", "What would you tell a friend?"),
    ("Q6: ACTION", "What can you do about it?"),
    ("SUMMARIZING", "Internal analysis (hidden)"),
    ("DELIVER_REFRAME", "Balanced thought from YOUR words"),
    ("RATE_NEW + RERATE OLD", "Measure belief shift"),
    ("RERATE_EMOTION", "Measure emotion change"),
    ("ACTION_PLAN", "Homework for the week"),
    ("COMPLETE", "Celebrate progress"),
]
y = 1.55
for s, desc in cog:
    add_card(slide, 5.0, y, 4.5, 0.28, CARD_BG)
    col = ACCENT if "SUMM" in s else BLUE
    add_text(slide, 5.1, y + 0.02, 2.3, 0.22, s, size=9, color=col, bold=True)
    add_text(slide, 7.4, y + 0.02, 2.0, 0.22, desc, size=9, color=LIGHT_GRAY)
    y += 0.31

# Right column - BA + Relapse + Post
x_r = 10.0
add_text(slide, x_r, 1.2, 3.0, 0.3, "BEHAVIORAL ACT. (3)", size=14, color=RED_ACCENT, bold=True)
ba = [
    ("BA_MONITORING", "Track daily activities"),
    ("BA_SCHEDULING", "Plan 1 small activity"),
    ("BA_GRADED_TASK", "Build graduated steps"),
]
y = 1.55
for s, desc in ba:
    add_card(slide, x_r, y, 3.1, 0.32, CARD_BG)
    add_text(slide, x_r + 0.1, y + 0.03, 1.7, 0.25, s, size=10, color=RED_ACCENT, bold=True)
    add_text(slide, x_r + 1.8, y + 0.03, 1.2, 0.25, desc, size=9, color=LIGHT_GRAY)
    y += 0.35

y += 0.15
add_text(slide, x_r, y, 3.0, 0.3, "RELAPSE (1)", size=14, color=GREEN, bold=True)
y += 0.3
add_card(slide, x_r, y, 3.1, 0.32, CARD_BG)
add_text(slide, x_r + 0.1, y + 0.03, 1.8, 0.25, "RELAPSE_PREVENTION", size=10, color=GREEN, bold=True)
add_text(slide, x_r + 2.0, y + 0.03, 1.0, 0.25, "Coping plan", size=9, color=LIGHT_GRAY)

y += 0.55
add_text(slide, x_r, y, 3.0, 0.3, "POST-SESSION (5)", size=14, color=RGBColor(0x26, 0x8D, 0x75), bold=True)
y += 0.3
post_items = [
    ("SCHEMA_CHECK", "Find core beliefs (Sess 4+)"),
    ("DRDT_OUTPUT", "Generate thought record"),
    ("SESSION_SUMMARY", "Patient states takeaways"),
    ("SESSION_FEEDBACK", "How did this session feel?"),
    ("SESSION_DONE", "Save & close"),
]
for s, desc in post_items:
    add_card(slide, x_r, y, 3.1, 0.32, CARD_BG)
    add_text(slide, x_r + 0.1, y + 0.03, 1.7, 0.25, s, size=10, color=RGBColor(0x26, 0x8D, 0x75), bold=True)
    add_text(slide, x_r + 1.8, y + 0.03, 1.2, 0.25, desc, size=9, color=LIGHT_GRAY)
    y += 0.35

# Bottom note
add_card(slide, 0.5, 6.7, 12.5, 0.55, CARD_BG)
add_text(slide, 0.7, 6.75, 12.0, 0.4,
         "Every session: BDI -> Pre-session -> Routing picks a path -> Core therapy -> Post-session closing",
         size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ==================== SLIDE 11: EXAMPLE CONVERSATION WALKTHROUGH ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "Example: Session Walkthrough (Moderate Path)", size=34, color=BLUE, bold=True)
add_text(slide, 0.8, 0.95, 11, 0.4,
         "Patient thinks \"I'm going to get fired\"  |  BDI = 22 (moderate)  |  Session 3",
         size=17, color=ORANGE)
add_divider(slide, 1.35)

walkthrough = [
    ("PRE-SESSION", ORANGE, [
        ("BDI_ASSESSMENT", "21 items scored -> BDI = 22 (moderate depression)"),
        ("BRIDGE", "\"How have things been since last time?\""),
        ("HOMEWORK_REVIEW", "Reviewed last week's thought record attempt"),
        ("AGENDA_SETTING", "\"Let's focus on the fear of getting fired\""),
        ("SEVERITY_ROUTING", "BDI = 22 -> routes to COGNITIVE path"),
    ]),
    ("AGENT 1: WARM QUESTIONER", BLUE, [
        ("VALIDATE", "\"That sounds really heavy. Thank you for sharing.\""),
        ("RATE_BELIEF + EMOTION", "Belief: 85%  |  Emotion: Anxiety at 90/100"),
        ("Q1: EVIDENCE", "For: bad feedback  |  Against: boss praised last project"),
        ("Q2: ALTERNATIVE", "\"Maybe he wants me to improve, not leave\""),
        ("Q3: WORST/BEST/REAL", "Fired / keeps job / gets another review in a month"),
        ("Q4: EFFECT", "\"Can't sleep, avoid my boss\""),
        ("Q5: FRIEND", "\"One bad review doesn't mean fired\""),
        ("Q6: ACTION", "\"Ask boss what to improve\""),
    ]),
    ("AGENT 2 + 3: REFRAME", GREEN, [
        ("SUMMARIZING", "Internal: Fortune Telling + Catastrophizing detected"),
        ("DELIVER_REFRAME", "\"One review doesn't define my career\" (from patient's own words)"),
        ("RERATE", "Belief: 85% -> 40%  |  Anxiety: 90 -> 50  (measurable improvement)"),
        ("ACTION_PLAN", "Homework: Ask boss for a quick improvement check-in"),
    ]),
    ("POST-SESSION", RGBColor(0x26, 0x8D, 0x75), [
        ("DRDT", "Auto-generated thought record worksheet"),
        ("SUMMARY + FEEDBACK", "Patient states takeaways in own words -> session saved"),
    ]),
]

y = 1.55
for phase, color, items in walkthrough:
    add_text(slide, 0.8, y, 4.0, 0.3, phase, size=14, color=color, bold=True)
    y += 0.32
    for state, desc in items:
        add_card(slide, 1.0, y, 11.5, 0.3, CARD_BG)
        add_text(slide, 1.15, y + 0.02, 2.5, 0.25, state, size=11, color=color, bold=True)
        add_text(slide, 3.7, y + 0.02, 8.5, 0.25, desc, size=11, color=LIGHT_GRAY)
        y += 0.32
    y += 0.1

# Result box
add_card(slide, 2.5, y + 0.05, 8.5, 0.45, RGBColor(0x1A, 0x3A, 0x1A))
add_text(slide, 2.7, y + 0.1, 8.0, 0.35,
         "Result: Belief dropped 85% -> 40%  |  Anxiety dropped 90 -> 50  |  Homework assigned",
         size=15, color=GREEN, bold=True, align=PP_ALIGN.CENTER)


# ==================== SLIDE 12: ALTERNATE PATHS (BA + RELAPSE) ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "Alternate Paths: Severe & Recovered Patients", size=34, color=BLUE, bold=True)
add_divider(slide, 1.0)

# Left: Behavioral Activation
add_card(slide, 0.7, 1.3, 6.0, 5.5, CARD_BG)
add_text(slide, 0.9, 1.4, 5.5, 0.35, "SEVERE (BDI >= 29): Behavioral Activation", size=18, color=RED_ACCENT, bold=True)
add_text(slide, 0.9, 1.85, 5.5, 0.35,
         "Too depressed for thought work. Focus on doing, not thinking.", size=14, color=LIGHT_GRAY)

ba_flow = [
    ("BA_MONITORING", "\"Walk me through a typical day\"", "Track what they do + rate Mastery & Pleasure"),
    ("BA_SCHEDULING", "\"Let's plan one tiny activity\"", "Specific: 5-min walk, tomorrow, 9am"),
    ("BA_GRADED_TASK", "\"Build up in small steps\"", "Step 1: 5 min -> Step 2: 10 min -> Step 3: 15 min"),
]
y = 2.4
for state, example, note in ba_flow:
    add_text(slide, 1.1, y, 2.2, 0.3, state, size=13, color=RED_ACCENT, bold=True)
    y += 0.3
    add_text(slide, 1.3, y, 5.0, 0.25, example, size=12, color=WHITE)
    y += 0.28
    add_text(slide, 1.3, y, 5.0, 0.25, note, size=11, color=LIGHT_GRAY)
    y += 0.45

add_card(slide, 1.0, y, 5.3, 0.55, RGBColor(0x3A, 0x1A, 0x1A))
add_text(slide, 1.2, y + 0.05, 5.0, 0.4,
         "Key: \"Activity precedes mood change\"\nSmall wins > big goals  (Beck et al., 1979)",
         size=12, color=RED_ACCENT)

# Right: Relapse Prevention
add_card(slide, 7.0, 1.3, 5.8, 5.5, CARD_BG)
add_text(slide, 7.2, 1.4, 5.3, 0.35, "RECOVERED (BDI < 14, 3+ sessions):", size=18, color=GREEN, bold=True)
add_text(slide, 7.2, 1.75, 5.3, 0.35, "Relapse Prevention", size=18, color=GREEN, bold=True)
add_text(slide, 7.2, 2.2, 5.3, 0.35,
         "Patient is getting better - prepare for independence.", size=14, color=LIGHT_GRAY)

relapse_items = [
    "1. Celebrate progress (\"BDI 35 -> 11, that's huge!\")",
    "2. Identify triggers (\"What might bring old patterns back?\")",
    "3. Build coping plan:",
    "     Warning signs: skipping meals, avoiding people",
    "     Strategies: thought record, walk, call a friend",
    "     Support contacts: friend, hotline",
    "4. Space out sessions (weekly -> monthly check-ins)",
    "5. Normalize setbacks (\"They don't erase progress\")",
]
y = 2.7
for item in relapse_items:
    add_text(slide, 7.4, y, 5.0, 0.3, item, size=13, color=LIGHT_GRAY)
    y += 0.35

add_card(slide, 7.3, y + 0.15, 5.2, 0.4, RGBColor(0x1A, 0x3A, 0x1A))
add_text(slide, 7.5, y + 0.2, 4.8, 0.3,
         "Goal: Patient can manage on their own",
         size=13, color=GREEN, bold=True)


# ==================== SLIDE 13: SAFETY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "Safety & Crisis Detection", size=36, color=BLUE, bold=True)
add_divider(slide, 1.1)

add_card(slide, 0.8, 1.5, 11.7, 2.5, CARD_BG)
add_text(slide, 1.1, 1.6, 11, 0.4, "Every message is screened for crisis signals", size=22, color=RED_ACCENT, bold=True)
safety_items = [
    "  BDI Item 9 (Suicidal Thoughts): Score >= 2 triggers immediate crisis response",
    "  Keyword-based crisis detection on every single chat message",
    "  Crisis events are flagged in database for admin review",
    "  Immediate crisis resources returned (hotlines, emergency contacts)",
    "  User name + email stored for follow-up by human professionals",
]
add_bullet_slide(slide, safety_items, top=2.15, spacing=0.45, size=18, color=LIGHT_GRAY)

add_card(slide, 0.8, 4.5, 11.7, 2.2, CARD_BG)
add_text(slide, 1.1, 4.6, 11, 0.4, "Quality Monitoring", size=22, color=GREEN, bold=True)
quality_items = [
    "  Supervisor agent (8B model) evaluates every therapist response",
    "  Checks: Collaboration, Guided Discovery, Focus, Warmth, Pacing",
    "  Flags safety concerns, fabricated techniques, state violations",
    "  Admin dashboard for monitoring all sessions and crisis flags",
]
add_bullet_slide(slide, quality_items, top=5.15, spacing=0.45, size=18, color=LIGHT_GRAY)

# ==================== SLIDE 14: REFERENCES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "Academic References", size=36, color=BLUE, bold=True)
add_divider(slide, 1.1)

refs = [
    "Beck, A.T., Rush, A.J., Shaw, B.F., & Emery, G. (1979).",
    "    Cognitive Therapy of Depression. Guilford Press.",
    "",
    "Beck, A.T., Steer, R.A., & Brown, G.K. (1996).",
    "    Manual for the Beck Depression Inventory-II. Psychological Corporation.",
    "",
    "Beck, J.S. (1995). Cognitive Behavior Therapy: Basics and Beyond. Guilford Press.",
    "",
    "Beck, J.S. (2020). Cognitive Behavior Therapy: Basics and Beyond, 3rd Edition.",
    "    Guilford Press. [Thought Record Worksheet - implemented in prompts.py]",
    "",
    "Beck, A.T. (1976). Cognitive Therapy and the Emotional Disorders. Penguin.",
]

y = 1.5
for ref in refs:
    if ref == "":
        y += 0.15
    else:
        color = LIGHT_GRAY if ref.startswith("    ") else WHITE
        bold = not ref.startswith("    ") and ref != ""
        add_text(slide, 1.0, y, 11, 0.4, ref, size=17, color=color, bold=bold)
        y += 0.38

add_card(slide, 0.8, y + 0.3, 11.7, 1.4, CARD_BG)
add_text(slide, 1.1, y + 0.4, 11, 0.4, "What this means:", size=20, color=ORANGE, bold=True)
add_text(slide, 1.1, y + 0.85, 11, 0.6,
         "Every technique in our app maps directly to a published Beck protocol.\n"
         "This is not \"inspired by CBT\" - it is a structured implementation of Beck's specific methods.",
         size=18, color=WHITE)

# ==================== SLIDE 15: CLOSING ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0, 2.2, 13.333, 1.0, "CBT Companion", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0, 3.3, 13.333, 0.6, "Beck's CBT Protocol, Implemented as a 32-State AI System",
         size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_divider(slide, 4.1)
add_text(slide, 0, 4.5, 13.333, 0.5, "BDI-II Assessment  |  6 Socratic Questions  |  Severity Routing",
         size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 0, 5.0, 13.333, 0.5, "Behavioral Activation  |  Schema Work  |  Thought Records",
         size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 0, 5.5, 13.333, 0.5, "Crisis Detection  |  Relapse Prevention",
         size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 0, 6.3, 13.333, 0.5, "Thank You", size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# Save
output_path = "CBT_Companion_Presentation_v3.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
